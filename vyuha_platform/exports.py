"""Turn one run into something you can hand to somebody.

Three formats, all rendered from the same ``Insights`` the dashboard uses, so a
deck, a PDF and an email can never disagree with the numbers on screen:

    to_pdf   — a printable one-pager for a meeting or a bank
    to_pptx  — a short deck for presenting to the owner
    draft_email — subject + body, ready to send or paste

Note on the rupee sign: the PDF core fonts (Helvetica) have no glyph for ₹ and
render it as a black box, so the PDF uses ``Rs.`` throughout. PowerPoint and
email use the real symbol. Same numbers, different alphabet — this is exactly
the split that ``vyuha.fmt`` exists to make explicit.
"""

from __future__ import annotations

import smtplib
from email.message import EmailMessage
from pathlib import Path

from vyuha import fmt
from vyuha.analyze import Insights

from . import channels

INK = "#111319"
ACCENT = "#7c5cff"
CRIT = "#c0392b"
WARN = "#b7791f"


def rs(v) -> str:
    """PDF-safe rupees: Helvetica has no glyph for the rupee sign."""
    return fmt.rupees(v or 0, symbol="Rs.")


def rs_short(v) -> str:
    return fmt.rupees_short(v or 0, symbol="Rs.")


# ------------------------------------------------------------------------ PDF

def to_pdf(insights: Insights, client: str, out: Path) -> Path:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_LEFT
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.platypus import (Paragraph, SimpleDocTemplate, Spacer, Table,
                                    TableStyle)

    out = Path(out)
    out.parent.mkdir(parents=True, exist_ok=True)
    base = getSampleStyleSheet()

    h1 = ParagraphStyle("h1", parent=base["Title"], fontSize=22, leading=26,
                        alignment=TA_LEFT, textColor=colors.HexColor(INK), spaceAfter=2)
    sub = ParagraphStyle("sub", parent=base["Normal"], fontSize=9.5, leading=13,
                         textColor=colors.HexColor("#6b7280"), spaceAfter=14)
    h2 = ParagraphStyle("h2", parent=base["Heading2"], fontSize=12, leading=15,
                        textColor=colors.HexColor(ACCENT), spaceBefore=16, spaceAfter=6)
    body = ParagraphStyle("body", parent=base["Normal"], fontSize=9.5, leading=13.5)

    doc = SimpleDocTemplate(str(out), pagesize=A4,
                            leftMargin=18 * mm, rightMargin=18 * mm,
                            topMargin=16 * mm, bottomMargin=16 * mm,
                            title=f"Vyuha brief - {client}", author="Vyuha")
    flow = [Paragraph(client or insights.source, h1),
            Paragraph(f"Vyuha operations brief &nbsp;·&nbsp; generated "
                      f"{insights.generated_at.strftime('%d %b %Y')}", sub)]

    # headline numbers
    s, st, r = insights.sales, insights.stock, insights.receivables
    cells = [["Revenue", "Orders", "Stock value", "Outstanding"],
             [rs_short(s.get("revenue")), str(s.get("orders", 0) or 0),
              rs_short(st.get("value")), rs_short(r.get("total"))]]
    tbl = Table(cells, colWidths=[43 * mm] * 4)
    tbl.setStyle(TableStyle([
        ("FONT", (0, 0), (-1, 0), "Helvetica", 7.5),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#6b7280")),
        ("FONT", (0, 1), (-1, 1), "Helvetica-Bold", 15),
        ("TEXTCOLOR", (0, 1), (-1, 1), colors.HexColor(INK)),
        ("BOTTOMPADDING", (0, 0), (-1, 0), 3),
        ("TOPPADDING", (0, 1), (-1, 1), 1),
        ("LINEBELOW", (0, 1), (-1, 1), 0.6, colors.HexColor("#e5e7eb")),
    ]))
    flow += [tbl, Spacer(1, 4)]

    # alerts
    flow.append(Paragraph("What needs attention", h2))
    alerts = channels.ordered(insights)
    if not alerts:
        flow.append(Paragraph("Nothing flagged on this file.", body))
    for a in alerts:
        colour = CRIT if a.severity == "critical" else WARN
        flow.append(Paragraph(
            f'<font color="{colour}"><b>{a.severity.upper()}</b></font> &nbsp; <b>'
            f'{_esc(a.title)}</b>', body))
        flow.append(Paragraph(_esc(a.detail).replace("₹", "Rs."), body))
        if a.entities:
            shown = ", ".join(_esc(e) for e in a.entities[:8]).replace("₹", "Rs.")
            more = f" (+{len(a.entities) - 8} more)" if len(a.entities) > 8 else ""
            flow.append(Paragraph(
                f'<font color="#6b7280">{shown}{more}</font>', body))
        flow.append(Spacer(1, 7))

    # top customers
    top = s.get("top_parties") or []
    if top:
        flow.append(Paragraph("Where the revenue comes from", h2))
        rows = [["Customer", "Revenue", "Share"]]
        rows += [[_esc(p["label"]), rs(p["amount"]), f"{p['share']:.0%}"] for p in top[:8]]
        t2 = Table(rows, colWidths=[95 * mm, 45 * mm, 32 * mm])
        t2.setStyle(TableStyle([
            ("FONT", (0, 0), (-1, 0), "Helvetica-Bold", 8),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#6b7280")),
            ("FONT", (0, 1), (-1, -1), "Helvetica", 9),
            ("ALIGN", (1, 0), (-1, -1), "RIGHT"),
            ("LINEBELOW", (0, 0), (-1, 0), 0.5, colors.HexColor("#e5e7eb")),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#fafafa")]),
            ("TOPPADDING", (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ]))
        flow.append(t2)

    flow += [Spacer(1, 14),
             Paragraph('<font color="#9ca3af" size="7.5">Produced by Vyuha from the file '
                       'the client supplied. No data was re-keyed.</font>', body)]
    doc.build(flow)
    return out


def _esc(text) -> str:
    return (str(text).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"))


# ----------------------------------------------------------------------- PPTX

def to_pptx(insights: Insights, client: str, out: Path) -> Path:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    out = Path(out)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    dark = RGBColor(0x0B, 0x0C, 0x10)
    white = RGBColor(0xF4, 0xF6, 0xFA)
    grey = RGBColor(0x9A, 0xA3, 0xB4)
    violet = RGBColor(0x7C, 0x5C, 0xFF)
    red = RGBColor(0xFB, 0x5F, 0x6D)
    amber = RGBColor(0xFB, 0xBF, 0x24)

    def slide(fill=dark):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = fill
        return s

    def text(s, x, y, w, h, body, size=18, colour=white, bold=False):
        box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(body)
        p.font.size, p.font.bold, p.font.color.rgb = Pt(size), bold, colour
        return tf

    # 1 — title
    s1 = slide()
    text(s1, 0.9, 2.4, 11.5, 1.2, client or insights.source, 46, white, True)
    text(s1, 0.9, 3.5, 11.5, 0.6, "Operations brief", 22, violet, True)
    text(s1, 0.9, 4.15, 11.5, 0.5,
         insights.generated_at.strftime("Generated %d %B %Y  ·  Vyuha"), 13, grey)

    # 2 — headline numbers
    s2 = slide()
    text(s2, 0.9, 0.6, 11.5, 0.6, "Where the business stands", 30, white, True)
    stats = [("Revenue", fmt.rupees_short(insights.sales.get("revenue") or 0)),
             ("Orders", str(insights.sales.get("orders", 0) or 0)),
             ("Stock value", fmt.rupees_short(insights.stock.get("value") or 0)),
             ("Outstanding", fmt.rupees_short(insights.receivables.get("total") or 0))]
    for i, (k, v) in enumerate(stats):
        x = 0.9 + i * 3.05
        text(s2, x, 2.3, 2.9, 0.4, k.upper(), 11, grey, True)
        text(s2, x, 2.75, 2.9, 1.0, v, 34, white, True)

    # 3 — alerts
    s3 = slide()
    text(s3, 0.9, 0.6, 11.5, 0.6, "What needs attention", 30, white, True)
    y = 1.7
    for a in channels.ordered(insights)[:5]:
        colour = red if a.severity == "critical" else amber
        text(s3, 0.9, y, 11.5, 0.45, a.title, 19, colour, True)
        text(s3, 0.9, y + 0.42, 11.5, 0.5, a.detail, 12.5, grey)
        y += 1.12
    if not insights.alerts:
        text(s3, 0.9, 1.9, 11.5, 0.6, "Nothing flagged on this file.", 18, grey)

    # 4 — concentration
    top = insights.sales.get("top_parties") or []
    if top:
        s4 = slide()
        text(s4, 0.9, 0.6, 11.5, 0.6, "Where the revenue comes from", 30, white, True)
        y = 1.7
        for p in top[:7]:
            text(s4, 0.9, y, 7.0, 0.4, p["label"], 15, white)
            text(s4, 8.2, y, 2.2, 0.4, fmt.rupees_short(p["amount"]), 15, white, True)
            text(s4, 10.8, y, 1.6, 0.4, f"{p['share']:.0%}", 15, violet, True)
            y += 0.62
        share = insights.sales.get("top3_share") or 0
        text(s4, 0.9, y + 0.3, 11.5, 0.5,
             f"Top 3 customers are {share:.0%} of revenue.", 13, grey)

    prs.save(str(out))
    return out


# ---------------------------------------------------------------------- email

def draft_email(insights: Insights, client: str, contact: str = "") -> tuple[str, str]:
    """Return (subject, body) — a framed email, not a raw metric dump."""
    when = insights.generated_at.strftime("%d %b %Y")
    alerts = channels.ordered(insights)
    critical = [a for a in alerts if a.severity == "critical"]

    if critical:
        subject = f"{client}: {len(critical)} thing(s) need attention - {when}"
    elif alerts:
        subject = f"{client}: {len(alerts)} item(s) to review - {when}"
    else:
        subject = f"{client}: everything looks steady - {when}"

    greeting = f"Hi {contact.split()[0]}," if contact.strip() else "Hi,"
    lines = [greeting, "",
             "We ran your latest file through Vyuha. Here is what stood out:", ""]
    for a in alerts:
        lines.append(f"* {a.title}")
        lines.append(f"  {a.detail}")
        if a.entities:
            lines.append(f"  {', '.join(a.entities[:6])}"
                         + (f" and {len(a.entities) - 6} more" if len(a.entities) > 6 else ""))
        lines.append("")
    if not alerts:
        lines += ["Nothing needs attention on this file.", ""]

    s, st, r = insights.sales, insights.stock, insights.receivables
    lines += ["---",
              f"Revenue {channels.money(s.get('revenue'))}   "
              f"Orders {s.get('orders', 0) or 0}",
              f"Stock {channels.money(st.get('value'))}   "
              f"Outstanding {channels.money(r.get('total'))}",
              "",
              "The full dashboard is attached. Nothing in your file was changed.",
              "", "- Vyuha"]
    return subject, "\n".join(lines)


def send_email(settings, to: str, subject: str, body: str,
               attachments: list[Path] | None = None) -> tuple[bool, str]:
    """Send via SMTP. Returns (ok, detail) — never raises."""
    if not settings.email_live:
        return False, "No SMTP server configured — use the draft to send it yourself."
    if not to.strip():
        return False, "No email address on file for this client."

    msg = EmailMessage()
    msg["From"] = settings.smtp_from
    msg["To"] = to
    msg["Subject"] = subject
    msg.set_content(body)

    for path in attachments or []:
        path = Path(path)
        if not path.exists():
            continue
        sub = {"html": ("text", "html"), "pdf": ("application", "pdf")}.get(
            path.suffix.lstrip("."), ("application", "octet-stream"))
        msg.add_attachment(path.read_bytes(), maintype=sub[0], subtype=sub[1],
                           filename=path.name)

    try:
        with smtplib.SMTP(settings.smtp_host, settings.smtp_port, timeout=30) as smtp:
            smtp.starttls()
            if settings.smtp_user:
                smtp.login(settings.smtp_user, settings.smtp_password)
            smtp.send_message(msg)
    except (smtplib.SMTPException, OSError) as exc:
        return False, f"{type(exc).__name__}: {exc}"
    return True, f"Sent to {to}."
