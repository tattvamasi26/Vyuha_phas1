"""A deck, from a sentence.

``exports.to_pptx`` already renders a fixed five-slide brief off ``Insights``.
That is the right thing when the answer is always "show me the operations
brief", and the wrong thing when the promoter says "make me a case study for
the bank" — a different audience, a different argument, the same numbers.

So this module separates the two jobs that were previously one:

* **What to say** — an ``Outline``: a title and a list of slides. Claude writes
  it from the brief plus the same compact facts the agent reads, so a deck can
  never contain a number the business does not have.
* **How it looks** — ``to_pptx`` / ``to_pdf`` render an ``Outline`` and know
  nothing about where it came from.

That split is what makes the offline path honest rather than a degraded mode:
``_fallback()`` builds a real outline from the facts with no model at all, and
it renders through exactly the same code. With no API key you still get a deck;
what you lose is the argument being tailored to the brief.

Every generated figure is passed through as a **string already formatted by
Python**. The model is asked to select and arrange, never to calculate.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path

from . import llm
from .agent import _rs

#: What the deck is for. The audience changes the argument far more than the
#: numbers do, which is why this is a choice and not a checkbox.
KINDS = {
    "review": ("Business review", "For the owner — how the business is doing and what to fix"),
    "case-study": ("Case study", "For a prospect — what changed, with the numbers behind it"),
    "investor": ("Investor brief", "For funding — traction, unit economics, where it goes"),
    "bank": ("Bank / credit file", "For a lender — cash flow, receivables, ability to repay"),
}

OUTLINE_SCHEMA = {
    "type": "object",
    "additionalProperties": False,
    "required": ["title", "subtitle", "slides"],
    "properties": {
        "title": {"type": "string"},
        "subtitle": {"type": "string"},
        "slides": {
            "type": "array",
            "minItems": 4,
            "maxItems": 8,
            "items": {
                "type": "object",
                "additionalProperties": False,
                "required": ["heading", "kind", "bullets", "stats"],
                "properties": {
                    "heading": {"type": "string"},
                    "kind": {
                        "type": "string",
                        "enum": ["stats", "points", "split", "closing"],
                        "description": ("stats when the point IS the figures; "
                                        "split to compare two sets; closing for "
                                        "the last slide; points otherwise."),
                    },
                    "bullets": {
                        "type": "array", "maxItems": 5,
                        "items": {"type": "string"},
                    },
                    "stats": {
                        "type": "array", "maxItems": 4,
                        "items": {
                            "type": "object",
                            "additionalProperties": False,
                            "required": ["label", "value"],
                            "properties": {
                                "label": {"type": "string"},
                                "value": {"type": "string"},
                            },
                        },
                    },
                },
            },
        },
    },
}

SYSTEM = """You write short, factual business decks for Indian small businesses.

Hard rules:
1. Every number you use must appear in the supplied JSON facts. Never invent,
   never estimate, never extrapolate a trend that is not in the data.
2. Numbers in `stats.value` must be copied as strings exactly as they appear in
   the facts (they are already formatted). Do not recompute or reformat them.
3. If the facts do not support a slide the brief asks for, write that slide with
   a bullet saying plainly what data would be needed. Do not fake it.
4. Bullets are one line each, under 90 characters, no trailing full stop.
5. No filler slides. No "Thank you", no "Agenda", no "Questions?".
6. Write in plain business English an owner would use, not consultant language.
7. Choose `kind` by what the slide is doing:
   - "stats" when the point IS the numbers. Give 2-4 stats and at most 2 bullets.
   - "split" to set two things against each other (this year vs last, us vs them).
   - "closing" for the final slide: what to do about all of it.
   - "points" otherwise.
   A trend slide is added for you from the real monthly figures — do not write one."""


@dataclass
class Slide:
    heading: str
    bullets: list[str] = field(default_factory=list)
    stats: list[dict] = field(default_factory=list)
    #: cover | stats | chart | split | points | closing. The argument decides
    #: this, not the styling: a figure wants the whole screen, a trend wants a
    #: chart, a comparison wants two columns. Bullets are what you fall back to
    #: when none of those fit — not the default, which is what made the first
    #: version look like text on a rectangle.
    kind: str = "points"
    #: {"type": "bar"|"hbar", "series": [{"label","value","value2","display"}],
    #:  "legend": [(label, colour)]} — drawn as inline SVG from real numbers.
    chart: dict | None = None


@dataclass
class Outline:
    title: str
    subtitle: str
    slides: list[Slide] = field(default_factory=list)
    #: "claude" | "cache" | "built" — shown before download, because a deck
    #: written to a brief and a deck assembled from a template are not the
    #: same thing and the operator is about to put their name on it.
    source: str = "built"
    note: str = ""

    @property
    def label(self) -> str:
        return {"claude": "Written by Claude from your brief",
                "cache": "Written by Claude (from cache)",
                "built": "Assembled from your numbers"}.get(self.source, self.source)


# ----------------------------------------------------------------- the charts

def _trend_slide(f: dict) -> "Slide | None":
    """Revenue and profit by month, from the numbers finance.py computed."""
    months = f.get("monthly_trend") or []
    if len(months) < 3:
        return None
    series = [{"label": m["month"].split()[0], "value": m["revenue"],
               "value2": m["profit"]} for m in months[-9:]]
    latest, first = months[-1], months[0]
    change = ((latest["revenue"] - first["revenue"]) / first["revenue"] * 100
              if first["revenue"] else 0)
    return Slide(
        heading="Month by month",
        kind="chart",
        chart={"type": "bar", "series": series,
               "legend": [("Revenue", "#F0A93C"), ("Profit", "#3FBFA4")]},
        bullets=[f"Latest month: {_rs(latest['revenue'])} revenue, "
                 f"{_rs(latest['profit'])} profit",
                 f"{'Up' if change >= 0 else 'Down'} {abs(change):.0f}% "
                 f"against {first['month']}"],
    )


def _customers_slide(f: dict) -> "Slide | None":
    """Who the revenue comes from — the slide that makes concentration visible."""
    rows = f.get("top_customers") or []
    if len(rows) < 2:
        return None
    risk = f.get("customer_concentration_risk", "spread")
    note = {"high": "One customer leaving would hurt badly.",
            "watch": "Worth widening the base.",
            "spread": "Well spread — no single point of failure."}[risk]
    return Slide(
        heading="Where the revenue comes from",
        kind="chart",
        chart={"type": "hbar",
               "series": [{"label": c["party"], "value": c["amount"],
                           "display": f"{c['share_pct']:.0f}%"} for c in rows[:6]]},
        bullets=[f"Top customer is {rows[0]['share_pct']:.0f}% of revenue.", note],
    )


# --------------------------------------------------------------- the fallback

def _fallback(kind: str, brief: str, client, f: dict) -> Outline:
    """A real deck with no model involved. Same renderer, same numbers.

    Not a degraded mode: it picks slide kinds, builds the same charts, and reads
    the same statements. What it cannot do is tailor the argument to the brief,
    and it says so rather than pretending.
    """
    sales, stock, cash = f["sales"], f["stock"], f["money"]
    chase = f["followups"]
    title = {"case-study": f"{client.name} — case study",
             "investor": f"{client.name} — investor brief",
             "bank": f"{client.name} — credit file"}.get(
                 kind, f"{client.name} — business review")

    slides = [
        Slide("Where the business stands", kind="stats",
              stats=[{"label": "Earned", "value": _rs(sales["total_earned"])},
                     {"label": "Collected", "value": _rs(sales["cash_collected"])},
                     {"label": "Owed to you", "value": _rs(sales["still_owed_to_you"])},
                     {"label": "Stock value", "value": _rs(stock["stock_value"])}],
              bullets=[f"{sales['bills']} bills from {sales['customers']} customers",
                       f"Figures as of {f['as_of']}"]),
    ]

    trend = _trend_slide(f)
    if trend is not None:
        slides.append(trend)

    slides.append(Slide(
        "Cash, and what it is not", kind="split",
        stats=[{"label": "Net cash", "value": _rs(cash["net_cash"])},
               {"label": "Profit", "value": _rs(sales["profit_where_cost_known"])}],
        bullets=[f"{_rs(cash['cash_came_in'])} came in",
                 f"{_rs(cash['cash_went_out'])} went out",
                 f"{_rs(cash['still_to_collect'])} still to collect",
                 f"{_rs(cash['still_to_pay'])} still to pay"]))

    customers = _customers_slide(f)
    if customers is not None:
        slides.append(customers)

    attention = ([f"{len(stock['below_reorder'])} item(s) below reorder level"]
                 if stock["below_reorder"] else [])
    attention += ([f"{len(stock['out_of_stock'])} item(s) out of stock"]
                  if stock["out_of_stock"] else [])
    attention += ([f"{_rs(chase['money_to_chase'])} overdue from "
                   f"{len(chase['overdue_payments'])} customer(s)"]
                  if chase.get("overdue_payments") else [])
    attention += ([f"{len(stock['never_sold'])} item(s) have never sold"]
                  if stock["never_sold"] else [])
    slides.append(Slide(
        "What needs attention", kind="points",
        bullets=attention or ["Nothing is flagged — stock, cash and collections "
                              "are all clean"]))

    if sales["best_sellers"]:
        slides.append(Slide(
            "What sells", kind="chart",
            chart={"type": "hbar",
                   "series": [{"label": b["item"], "value": b["qty_sold"],
                               "display": f"{b['qty_sold']:g}"}
                              for b in sales["best_sellers"][:6]]},
            bullets=[f"{sales['best_sellers'][0]['item']} leads on volume."]))

    if f.get("branches", {}).get("branches"):
        slides.append(Slide(
            "By branch", kind="chart",
            chart={"type": "hbar",
                   "series": [{"label": b["name"], "value": b["revenue"]}
                              for b in f["branches"]["branches"][:6]]}))

    slides.append(Slide(
        "What to do about it", kind="closing",
        bullets=(attention[:3] or ["Keep going — nothing is flagged"])))

    return Outline(title=title,
                   subtitle=brief.strip()[:120] or KINDS.get(kind, ("Business review", ""))[0],
                   slides=slides, source="built",
                   note="Assembled from your numbers — add a Claude API key to "
                        "have it argued to your brief.")


# ------------------------------------------------------------------- building

def outline(brief: str, kind: str, client, f: dict, settings) -> Outline:
    """Decide what the deck says. Falls back to a built outline, never fails."""
    kind = kind if kind in KINDS else "review"
    audience = KINDS[kind][1]

    import json
    prompt = (
        f"Business: {client.name}\n"
        f"Deck type: {KINDS[kind][0]} — {audience}\n"
        f"What the promoter asked for: {brief.strip() or 'A general review of the business.'}\n\n"
        f"Facts (every number you use must come from here):\n"
        f"```json\n{json.dumps(f, indent=2, default=str)}\n```\n\n"
        f"Write 4 to 6 slides. Amounts are Indian rupees."
    )

    answer = llm.ask(prompt, settings, system=SYSTEM, schema=OUTLINE_SCHEMA)
    if not answer.ok or not answer.data:
        out = _fallback(kind, brief, client, f)
        if answer.error:
            out.note = answer.error + (f" {answer.needs_action}" if answer.needs_action else "")
        return out

    data = answer.data
    slides = [
        Slide(heading=str(s.get("heading", "")).strip(),
              kind=str(s.get("kind", "points")) or "points",
              bullets=[str(b).strip() for b in s.get("bullets", []) if str(b).strip()],
              stats=[{"label": str(x.get("label", "")), "value": str(x.get("value", ""))}
                     for x in s.get("stats", [])])
        for s in data.get("slides", [])
    ]
    slides = [s for s in slides if s.heading]
    if not slides:
        return _fallback(kind, brief, client, f)

    # The trend slide is built here, from the real series, and inserted after
    # the opener. Asking a model to hand back chart data invites a plotted
    # number that no statement agrees with.
    trend = _trend_slide(f)
    if trend is not None:
        slides.insert(min(1, len(slides)), trend)

    return Outline(title=str(data.get("title") or client.name),
                   subtitle=str(data.get("subtitle") or KINDS[kind][0]),
                   slides=slides,
                   source="cache" if answer.cached else "claude")


# ------------------------------------------------------------------ rendering

def to_pptx(outline: Outline, client_name: str, out: Path) -> Path:
    """Same visual language as exports.to_pptx — one product, one look."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    out = Path(out)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    dark = RGBColor(0x0B, 0x0C, 0x10)
    white = RGBColor(0xF4, 0xF6, 0xFA)
    grey = RGBColor(0x9A, 0xA3, 0xB4)
    violet = RGBColor(0x7C, 0x5C, 0xFF)
    amber = RGBColor(0xF0, 0xA9, 0x3C)

    def slide():
        s = prs.slides.add_slide(prs.slide_layouts[6])
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = dark
        return s

    def text(s, x, y, w, h, body, size=18, colour=white, bold=False):
        box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(body)
        p.font.size, p.font.bold, p.font.color.rgb = Pt(size), bold, colour
        return tf

    # title
    s = slide()
    text(s, 0.9, 2.4, 11.5, 1.3, outline.title, 44, white, True)
    text(s, 0.9, 3.7, 11.5, 0.7, outline.subtitle, 20, violet, True)
    text(s, 0.9, 4.5, 11.5, 0.5,
         datetime.now().strftime("%d %B %Y") + "  ·  Vyuha", 13, grey)

    def bar(s, x, y, w, h, colour=violet):
        """A drawn rectangle — python-pptx has charts, but a plain shape is
        fewer moving parts and looks the same at this size."""
        from pptx.enum.shapes import MSO_SHAPE
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(y), Inches(max(w, 0.04)),
                                   Inches(max(h, 0.04)))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colour
        shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    for sl in outline.slides:
        s = slide()
        text(s, 0.9, 0.6, 11.5, 0.8, sl.heading, 30, white, True)
        y = 1.8
        if sl.stats:
            width = min(3.05, 11.5 / max(len(sl.stats), 1))
            for i, stat in enumerate(sl.stats):
                x = 0.9 + i * width
                text(s, x, y, width - 0.15, 0.4, stat["label"].upper(), 11, grey, True)
                text(s, x, y + 0.42, width - 0.15, 0.9, stat["value"], 30, white, True)
            y += 1.7

        # A chart slide carries its argument in the chart. Without this it
        # exported as a heading on an empty rectangle -- exactly the "too basic"
        # the HTML deck was rebuilt to fix, reintroduced in the file people
        # actually open in PowerPoint.
        chart = getattr(sl, "chart", None)
        if chart and chart.get("series"):
            series = chart["series"][:9]
            peak = max((abs(float(p.get("value") or 0)) for p in series), default=1) or 1
            if chart.get("type") == "hbar":
                row_h, gap = 0.34, 0.12
                for point in series:
                    value = float(point.get("value") or 0)
                    text(s, 0.9, y, 2.9, row_h, str(point.get("label", ""))[:26], 13, grey)
                    bar(s, 3.9, y + 0.06, 7.0 * (abs(value) / peak), row_h - 0.12)
                    text(s, 11.0, y, 1.4, row_h,
                         str(point.get("display") or f"{value:,.0f}"), 12, white, True)
                    y += row_h + gap
            else:
                plot_h, base = 2.7, y + 2.9
                slot = 11.5 / max(len(series), 1)
                for i, point in enumerate(series):
                    value = float(point.get("value") or 0)
                    height = plot_h * (abs(value) / peak)
                    x = 0.9 + i * slot
                    bar(s, x, base - height, slot * 0.30, height)
                    if point.get("value2") is not None:
                        v2 = abs(float(point["value2"]))
                        h2 = plot_h * (v2 / peak)
                        bar(s, x + slot * 0.34, base - h2, slot * 0.30, h2, amber)
                    text(s, x, base + 0.06, slot * 0.9, 0.3,
                         str(point.get("label", ""))[:8], 11, grey)
                y = base + 0.5

        for bullet in sl.bullets:
            text(s, 0.9, y, 11.5, 0.5, "— " + bullet, 16, grey)
            y += 0.62

    prs.save(str(out))
    return out


def to_pdf(outline: Outline, client_name: str, out: Path) -> Path:
    """One page per slide. ``Rs.`` not ₹ — Helvetica has no rupee glyph."""
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.pdfgen import canvas as pdfcanvas

    out = Path(out)
    out.parent.mkdir(parents=True, exist_ok=True)
    width, height = landscape(A4)
    c = pdfcanvas.Canvas(str(out), pagesize=landscape(A4))

    ink = colors.HexColor("#111319")
    accent = colors.HexColor("#7c5cff")
    muted = colors.HexColor("#6b7280")

    def page_title(heading: str):
        c.setFillColor(ink)
        c.setFont("Helvetica-Bold", 24)
        c.drawString(48, height - 70, heading[:70])
        c.setStrokeColor(accent)
        c.setLineWidth(3)
        c.line(48, height - 84, 148, height - 84)

    # cover
    c.setFillColor(ink)
    c.setFont("Helvetica-Bold", 34)
    c.drawString(48, height / 2 + 20, outline.title[:60])
    c.setFillColor(accent)
    c.setFont("Helvetica-Bold", 16)
    c.drawString(48, height / 2 - 10, outline.subtitle[:90])
    c.setFillColor(muted)
    c.setFont("Helvetica", 11)
    c.drawString(48, height / 2 - 34, datetime.now().strftime("%d %B %Y") + "  -  Vyuha")
    c.showPage()

    for sl in outline.slides:
        page_title(sl.heading)
        y = height - 130
        if sl.stats:
            for i, stat in enumerate(sl.stats):
                x = 48 + i * 180
                c.setFillColor(muted)
                c.setFont("Helvetica-Bold", 9)
                c.drawString(x, y, stat["label"].upper()[:22])
                c.setFillColor(ink)
                c.setFont("Helvetica-Bold", 20)
                # The renderer must not smuggle in a glyph Helvetica lacks.
                c.drawString(x, y - 26, stat["value"].replace("₹", "Rs ")[:18])
            y -= 70
        c.setFont("Helvetica", 12)
        for bullet in sl.bullets:
            c.setFillColor(accent)
            c.drawString(48, y, "-")
            c.setFillColor(ink)
            c.drawString(64, y, bullet.replace("₹", "Rs ")[:110])
            y -= 22
            if y < 60:
                break
        c.showPage()

    c.save()
    return out
